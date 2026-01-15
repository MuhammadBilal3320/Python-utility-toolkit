from pptx import Presentation
from pptx.util import Inches, Pt
import os

# Create presentation object
prs = Presentation()

# Function to add a slide with title and detailed content
def add_slide(title_text, bullet_points):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title_text
    
    tf = slide.placeholders[1].text_frame
    tf.word_wrap = True
    tf.clear()
    
    for i, point in enumerate(bullet_points):
        if i == 0:
            tf.text = point
        else:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0

# --- Slide 1: Title Slide ---
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Sharpening Filters (Laplacian, High-Boost)"
slide.placeholders[1].text = (
    "Subject: Digital Image Processing\n"
    "Topic: Sharpening Filters\n\n"
    "Sharpening filters are image enhancement techniques used to emphasize "
    "fine details, edges, and intensity transitions. These filters improve "
    "the visual quality of images by highlighting high-frequency components."
)

# --- Slide 2: Introduction to Image Sharpening ---
add_slide("Introduction to Image Sharpening", [
    "Image sharpening is a process used to enhance edges and small details in an image.",
    "It focuses on emphasizing rapid intensity changes between neighboring pixels.",
    "Sharpening improves clarity and visual perception of important features.",
    "It is widely used in preprocessing stages of image analysis and computer vision.",
    "However, excessive sharpening can amplify noise and artifacts."
])

# --- Slide 3: Fundamentals of Sharpening Filters ---
add_slide("Fundamentals of Sharpening Filters", [
    "Sharpening filters operate by enhancing high-frequency components of an image.",
    "High frequencies correspond to edges, fine textures, and abrupt intensity changes.",
    "Derivative-based operations are commonly used for sharpening.",
    "First-order and second-order derivatives highlight different edge characteristics.",
    "Noise also contains high-frequency components, making noise control important."
])

# --- Slide 4: Laplacian Filter ---
add_slide("Laplacian Filter", [
    "The Laplacian filter is a second-order derivative operator.",
    "It measures the rate of change of the gradient in an image.",
    "The Laplacian responds strongly to regions with rapid intensity variation.",
    "It is isotropic, meaning it responds equally in all directions.",
    "Common Laplacian masks include 4-neighbor and 8-neighbor configurations."
])

# --- Slide 5: Laplacian-Based Image Sharpening ---
add_slide("Laplacian-Based Image Sharpening", [
    "Laplacian sharpening involves adding or subtracting the Laplacian result from the original image.",
    "This process enhances edges and fine details significantly.",
    "The method highlights both positive and negative intensity transitions.",
    "It is highly sensitive to noise due to second-order differentiation.",
    "Smoothing filters are often applied before Laplacian sharpening to reduce noise effects."
])

# --- Slide 6: High-Boost Filtering ---
add_slide("High-Boost Filtering", [
    "High-boost filtering is an advanced form of unsharp masking.",
    "It enhances edges while retaining more original image information.",
    "A blurred version of the image is subtracted from the original image.",
    "The difference image is scaled by a boost factor greater than one.",
    "The boost factor controls the strength of sharpening."
])

# --- Slide 7: Comparison of Laplacian and High-Boost ---
add_slide("Comparison: Laplacian vs High-Boost Filtering", [
    "Laplacian sharpening is based on second-order derivatives.",
    "High-boost filtering is based on subtracting a low-pass filtered image.",
    "Laplacian provides strong edge enhancement but is noise sensitive.",
    "High-boost allows better control over enhancement using the boost factor.",
    "High-boost generally produces smoother and more visually pleasing results."
])

# --- Slide 8: Applications and Conclusion ---
add_slide("Applications and Conclusion", [
    "Sharpening filters are widely used in medical imaging to enhance anatomical details.",
    "They are applied in satellite and remote sensing images for feature enhancement.",
    "Used in document processing to improve text clarity.",
    "Essential in computer vision and pattern recognition tasks.",
    "Laplacian and high-boost filters are powerful tools when used carefully."
])

# Save presentation in current directory
file_name = "Sharpening_Filters_Laplacian_HighBoost.pptx"
prs.save(file_name)

current_dir = os.getcwd()
full_path = os.path.join(current_dir, file_name)

print("Presentation saved in current directory:")
print(full_path)
